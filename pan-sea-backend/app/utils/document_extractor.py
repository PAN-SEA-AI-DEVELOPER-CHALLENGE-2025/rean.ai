import os
from typing import List


class DocumentExtractor:
    """Improved extractor using better parsers with safety caps and basic layout preservation."""

    @staticmethod
    def extract_text_from_file(
        file_path: str,
        *,
        max_pages: int = 200,
        max_slides: int = 200,
        max_chars: int = 50000,
    ) -> str:
        ext = os.path.splitext(file_path)[-1].lower()
        try:
            if ext == ".pdf":
                # Prefer PyMuPDF (fitz) for layout-aware extraction; fallback to pdfplumber, then pypdf
                try:
                    import fitz  # type: ignore
                    texts: List[str] = []
                    with fitz.open(file_path) as doc:
                        for i, page in enumerate(doc):
                            if i >= max_pages:
                                break
                            try:
                                page_text = page.get_text("blocks")  # list of (x0,y0,x1,y1, text, block_no)
                                blocks = [b[4] for b in page_text if len(b) > 4]
                                joined = "\n".join(t.strip() for t in blocks if t and t.strip())
                                if joined:
                                    texts.append(joined)
                            except Exception:
                                continue
                    content = "\n\n".join(texts).strip()
                    if content:
                        return content[:max_chars]
                except Exception:
                    pass
                try:
                    import pdfplumber  # type: ignore
                    texts: List[str] = []
                    with pdfplumber.open(file_path) as pdf:
                        for i, page in enumerate(pdf.pages):
                            if i >= max_pages:
                                break
                            try:
                                page_text = page.extract_text() or ""
                                if page_text:
                                    texts.append(page_text)
                            except Exception:
                                continue
                    content = "\n\n".join(texts).strip()
                    if content:
                        return content[:max_chars]
                except Exception:
                    pass
                try:
                    from pypdf import PdfReader  # type: ignore
                    reader = PdfReader(file_path)
                    texts: List[str] = []
                    for i, page in enumerate(getattr(reader, "pages", [])):
                        if i >= max_pages:
                            break
                        try:
                            page_text = page.extract_text() or ""
                            if page_text:
                                texts.append(page_text)
                        except Exception:
                            continue
                    content = "\n\n".join(texts).strip()
                    return content[:max_chars]
                except Exception:
                    return ""

            if ext in {".pptx"}:
                try:
                    from pptx import Presentation  # type: ignore
                except Exception:  # pragma: no cover
                    return ""
                try:
                    prs = Presentation(file_path)
                    texts: List[str] = []
                    for i, slide in enumerate(prs.slides):
                        if i >= max_slides:
                            break
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, "text"):
                                    value = getattr(shape, "text") or ""
                                    if value:
                                        texts.append(value)
                            except Exception:
                                continue
                    content = "\n\n".join(texts).strip()
                    return content[:max_chars]
                except Exception:
                    return ""

            if ext in {".docx"}:
                try:
                    from docx import Document  # type: ignore
                except Exception:  # pragma: no cover
                    return ""
                try:
                    doc = Document(file_path)
                    paragraphs = [p.text for p in getattr(doc, "paragraphs", []) if getattr(p, "text", "").strip()]
                    content = "\n\n".join(paragraphs).strip()
                    return content[:max_chars]
                except Exception:
                    return ""

            if ext in {".txt"}:
                try:
                    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                        return f.read().strip()[:max_chars]
                except Exception:
                    return ""

            # Unsupported extension
            return ""
        except Exception:
            return ""

    @staticmethod
    def extract_text_from_files(file_paths: List[str], *, max_total_chars: int = 80000) -> str:
        texts: List[str] = []
        for path in file_paths or []:
            content = DocumentExtractor.extract_text_from_file(path)
            if content:
                filename = os.path.basename(path)
                texts.append(f"[Material: {filename}]\n{content}")
        combined = "\n\n".join(texts).strip()
        return combined[:max_total_chars]



